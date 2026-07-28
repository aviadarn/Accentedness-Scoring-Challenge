#!/usr/bin/env python3
"""Generate the editable take-home challenge presentation.

The deck intentionally uses only native PowerPoint shapes and text so every
element remains editable after generation.
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = Path(__file__).with_name("accentedness-scoring-challenge.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_HEAD = "Avenir Next"
FONT_BODY = "Avenir Next"
FONT_MONO = "Menlo"

NAVY = "081321"
NAVY_2 = "0D1C2D"
CARD = "12263A"
CARD_2 = "173047"
LINE = "28435C"
WHITE = "F7FAFC"
MUTED = "9EB0C3"
AQUA = "5FE3D2"
BLUE = "73B7FF"
CORAL = "FF9274"
GOLD = "FFC857"
GREEN = "A8E063"
RED = "FF6B6B"
INK = "142033"


def inch(value: float):
    return Inches(value)


def color(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(NAVY)
    return slide


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = CARD,
    line: str | None = None,
    radius: bool = True,
    line_width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color(line)
        shape.line.width = Pt(line_width)
    return shape


def add_outline_oval(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    line: str = AQUA,
    line_width: float = 2.0,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(w), inch(h))
    shape.fill.background()
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    line: str = LINE,
    width: float = 1.5,
):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inch(x1),
        inch(y1),
        inch(x2),
        inch(y2),
    )
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(width)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    fill: str = WHITE,
    bold: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color(fill)
    return box


def add_rich_text(
    slide,
    segments: list[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = inch(0.02)
    tf.margin_top = tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    for value, style in segments:
        run = p.add_run()
        run.text = value
        run.font.name = style.get("font", FONT_BODY)
        run.font.size = Pt(style.get("size", size))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = color(style.get("fill", WHITE))
    return box


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = CARD_2,
    text_color: str = AQUA,
    size: float = 10,
    line: str | None = None,
):
    add_rect(slide, x, y, w, 0.32, fill=fill, line=line, radius=True)
    add_text(
        slide,
        text,
        x + 0.05,
        y + 0.01,
        w - 0.10,
        0.29,
        size=size,
        fill=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_dot_bullet(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: str = AQUA,
    size: float = 14,
    h: float = 0.5,
):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y + 0.12), inch(0.10), inch(0.10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color(accent)
    dot.line.fill.background()
    add_text(slide, text, x + 0.20, y, w - 0.20, h, size=size, fill=WHITE)


def add_waveform(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    stroke: str = AQUA,
    bars: int = 28,
):
    for idx in range(bars):
        phase = (idx / max(1, bars - 1)) * math.pi * 3.4
        envelope = 0.25 + 0.70 * math.sin(math.pi * idx / max(1, bars - 1)) ** 2
        magnitude = (0.35 + 0.65 * abs(math.sin(phase))) * envelope
        bar_h = max(0.05, h * magnitude)
        bar_w = w / bars * 0.42
        bar_x = x + idx * (w / bars) + (w / bars - bar_w) / 2
        bar_y = y + (h - bar_h) / 2
        add_rect(slide, bar_x, bar_y, bar_w, bar_h, fill=stroke, radius=True)


def add_header(slide, title: str, number: int, kicker: str | None = None):
    add_text(
        slide,
        kicker or "PHONE-LEVEL ACCENTEDNESS • TAKE-HOME",
        0.68,
        0.35,
        6.0,
        0.26,
        size=9,
        fill=AQUA,
        bold=True,
    )
    add_text(slide, title, 0.68, 0.77, 11.7, 0.72, size=28, fill=WHITE, bold=True, font=FONT_HEAD)
    add_text(
        slide,
        f"{number:02d}",
        12.15,
        0.36,
        0.48,
        0.25,
        size=9,
        fill=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_line(slide, 0.68, 1.55, 12.65, 1.55, line=LINE, width=1.0)


def add_footer(slide, source: str):
    add_line(slide, 0.68, 7.09, 12.65, 7.09, line=LINE, width=0.7)
    add_text(slide, source, 0.68, 7.15, 9.8, 0.18, size=6.5, fill=MUTED)
    add_text(
        slide,
        "github.com/aviadarn/Accentedness-Scoring-Challenge",
        9.25,
        7.15,
        3.40,
        0.18,
        size=6.5,
        fill=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def metric_card(slide, value: str, label: str, x: float, y: float, w: float, *, accent: str = AQUA):
    add_rect(slide, x, y, w, 1.10, fill=CARD, line=LINE)
    add_text(slide, value, x + 0.18, y + 0.15, w - 0.36, 0.48, size=24, fill=accent, bold=True)
    add_text(slide, label, x + 0.18, y + 0.70, w - 0.36, 0.24, size=10, fill=MUTED, bold=True)


def score_chip(slide, phone: str, score: str, x: float, y: float, *, accent: str):
    add_rect(slide, x, y, 0.73, 0.88, fill=CARD_2, line=accent, radius=True, line_width=1.4)
    add_text(slide, phone, x, y + 0.08, 0.73, 0.30, size=13, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, score, x, y + 0.48, 0.73, 0.23, size=10, fill=accent, bold=True, align=PP_ALIGN.CENTER)


def slide_01(prs: Presentation):
    slide = blank_slide(prs)
    add_outline_oval(slide, 9.35, -0.35, 4.55, 4.55, line=LINE, line_width=1.4)
    add_outline_oval(slide, 10.20, 4.95, 3.40, 3.40, line=AQUA, line_width=1.1)
    add_text(slide, "ML TAKE-HOME • 2026", 0.72, 0.58, 3.8, 0.28, size=10, fill=AQUA, bold=True)
    add_text(
        slide,
        "Phone-level\naccentedness scoring",
        0.72,
        1.20,
        7.4,
        1.72,
        size=38,
        fill=WHITE,
        bold=True,
        font=FONT_HEAD,
    )
    add_rect(slide, 0.72, 3.16, 1.15, 0.08, fill=AQUA, radius=False)
    add_text(
        slide,
        "From speech and expected phonemes to an interpretable\n0–100 score for every sound.",
        0.72,
        3.46,
        6.9,
        0.88,
        size=19,
        fill=MUTED,
    )
    add_pill(slide, "AUDIO + PHONES", 0.72, 4.62, 1.62)
    add_pill(slide, "ORDINAL MODEL", 2.50, 4.62, 1.62, text_color=BLUE)
    add_pill(slide, "GRADIO DEMO", 4.28, 4.62, 1.55, text_color=CORAL)
    add_text(slide, "Aviad Arnias", 0.72, 6.18, 3.0, 0.36, size=15, fill=WHITE, bold=True)
    add_text(slide, "Modeling • evaluation • product demo", 0.72, 6.57, 4.6, 0.26, size=10, fill=MUTED)

    add_rect(slide, 8.48, 0.78, 4.12, 5.85, fill=NAVY_2, line=LINE)
    add_text(slide, "ILLUSTRATIVE OUTPUT", 8.83, 1.13, 2.55, 0.24, size=9, fill=MUTED, bold=True)
    add_waveform(slide, 8.84, 1.65, 3.35, 1.20, stroke=AQUA)
    add_text(slide, "hello", 8.84, 3.10, 2.0, 0.40, size=20, fill=WHITE, bold=True)
    add_text(slide, "expected phones", 10.65, 3.20, 1.52, 0.22, size=8, fill=MUTED, align=PP_ALIGN.RIGHT)
    score_chip(slide, "h", "92", 8.84, 3.82, accent=GREEN)
    score_chip(slide, "ə", "45", 9.73, 3.82, accent=GOLD)
    score_chip(slide, "l", "89", 10.62, 3.82, accent=GREEN)
    score_chip(slide, "oʊ", "12", 11.51, 3.82, accent=CORAL)
    add_rect(slide, 8.84, 5.18, 3.35, 0.82, fill=CARD, line=LINE)
    add_text(slide, "Score = 50 × [P(Y≥1) + P(Y≥2)]", 9.04, 5.38, 2.95, 0.32, size=12, fill=BLUE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "Source: challenge brief example; scores shown are illustrative.", 8.84, 6.25, 3.35, 0.20, size=6.5, fill=MUTED, align=PP_ALIGN.CENTER)
    return slide


def slide_02(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "The task: localize pronunciation feedback", 2)
    cards = [
        (0.72, "01", "Speech", "Record or upload\na learner utterance", AQUA),
        (4.46, "02", "Expected phones", "The intended phone\nsequence—no timestamps", BLUE),
        (8.20, "03", "Phone scores", "One continuous 0–100\nscore per expected phone", CORAL),
    ]
    for x, num, title, body, accent in cards:
        add_rect(slide, x, 2.00, 3.08, 2.35, fill=CARD, line=LINE)
        add_pill(slide, num, x + 0.22, 2.20, 0.54, fill=accent, text_color=NAVY, size=10)
        add_text(slide, title, x + 0.22, 2.77, 2.55, 0.40, size=20, fill=WHITE, bold=True)
        add_text(slide, body, x + 0.22, 3.35, 2.58, 0.72, size=13, fill=MUTED)
    for x in (3.94, 7.68):
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, inch(x), inch(2.82), inch(0.34), inch(0.70))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color(LINE)
        arrow.line.fill.background()

    add_text(slide, "What makes it difficult", 0.72, 4.83, 3.3, 0.40, size=17, fill=WHITE, bold=True)
    add_dot_bullet(slide, "Labels are ordered (0 < 1 < 2), but output must be continuous.", 0.75, 5.35, 3.80, accent=AQUA, size=12, h=0.58)
    add_dot_bullet(slide, "Phone boundaries are absent and must be inferred from audio.", 4.56, 5.35, 3.70, accent=BLUE, size=12, h=0.58)
    add_dot_bullet(slide, "The majority label dominates, so raw accuracy is misleading.", 8.28, 5.35, 3.62, accent=CORAL, size=12, h=0.58)
    add_rect(slide, 0.72, 6.23, 11.90, 0.55, fill=CARD_2, line=LINE)
    add_text(slide, "Product question", 0.94, 6.38, 1.23, 0.20, size=9, fill=AQUA, bold=True)
    add_text(slide, "Where did this realization sound less American-like—and what should the learner practice next?", 2.28, 6.31, 9.98, 0.29, size=13, fill=WHITE, bold=True)
    add_footer(slide, "Sources: data/phone-scoring-ml-challenge.md; submission/WRITEUP.md")
    return slide


def slide_03(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "The data is small—and heavily imbalanced", 3)
    add_text(slide, "TRAIN LABEL DISTRIBUTION", 0.82, 1.93, 4.2, 0.25, size=9, fill=MUTED, bold=True)
    chart_x, chart_y, chart_w, chart_h = 0.92, 2.35, 5.65, 3.18
    add_line(slide, chart_x, chart_y + chart_h, chart_x + chart_w, chart_y + chart_h, line=LINE)
    values = [("0", 12.23, "heavily accented", CORAL), ("1", 7.88, "understandable", GOLD), ("2", 79.89, "native-like", AQUA)]
    max_value = 85.0
    for idx, (label, value, desc, accent) in enumerate(values):
        x = chart_x + 0.52 + idx * 1.72
        h = chart_h * value / max_value
        add_rect(slide, x, chart_y + chart_h - h, 0.92, h, fill=accent, radius=True)
        add_text(slide, f"{value:.2f}%", x - 0.12, chart_y + chart_h - h - 0.37, 1.16, 0.25, size=12, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.25, chart_y + chart_h + 0.14, 0.42, 0.26, size=14, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x - 0.30, chart_y + chart_h + 0.46, 1.52, 0.28, size=8, fill=MUTED, align=PP_ALIGN.CENTER)

    metric_card(slide, "2,799", "TRAIN UTTERANCES", 7.02, 1.96, 2.45, accent=AQUA)
    metric_card(slide, "100", "VALIDATION UTTERANCES", 9.73, 1.96, 2.45, accent=BLUE)
    metric_card(slide, "90,239", "LABELED PHONES TOTAL", 7.02, 3.34, 2.45, accent=GOLD)
    metric_card(slide, "44", "PHONE TYPES", 9.73, 3.34, 2.45, accent=CORAL)
    add_rect(slide, 7.02, 4.74, 5.16, 1.12, fill=CARD_2, line=CORAL, line_width=1.3)
    add_text(slide, "79.47%", 7.28, 4.94, 1.45, 0.40, size=25, fill=CORAL, bold=True)
    add_text(slide, "validation accuracy by always predicting “native-like”", 8.75, 4.91, 3.08, 0.42, size=11, fill=WHITE, bold=True)
    add_text(slide, "Yet recall for labels 0 and 1 is zero.", 8.75, 5.43, 3.08, 0.20, size=9, fill=MUTED)

    add_pill(slide, "NO PHONE TIMESTAMPS", 0.82, 6.34, 1.86, text_color=MUTED)
    add_pill(slide, "NO SPEAKER IDs", 2.84, 6.34, 1.42, text_color=MUTED)
    add_pill(slide, "NO RATER AGREEMENT", 4.42, 6.34, 1.82, text_color=MUTED)
    add_text(slide, "These gaps shape both modeling and evaluation.", 6.56, 6.38, 5.60, 0.22, size=11, fill=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, "Sources: data/README.md; submission/WRITEUP.md; submission/model/metrics.json")
    return slide


def slide_04(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Architecture: align first, then score", 4)
    stages = [
        ("01", "Audio", "16 kHz\nwaveform", AQUA),
        ("02", "Whisper-tiny", "Frozen 384-d\nencoder states", BLUE),
        ("03", "CTC + Viterbi", "45-way head;\nconstrained spans", GOLD),
        ("04", "Phone features", "mean + std +\n4 diagnostics", CORAL),
        ("05", "BiGRU", "phone embedding;\n2 layers", BLUE),
        ("06", "Ordinal head", "q₁, q₂ →\nscore 0–100", AQUA),
    ]
    x0, gap, width = 0.66, 0.17, 1.84
    for idx, (num, title, body, accent) in enumerate(stages):
        x = x0 + idx * (width + gap)
        add_rect(slide, x, 2.05, width, 2.60, fill=CARD, line=LINE)
        add_pill(slide, num, x + 0.17, 2.24, 0.48, fill=accent, text_color=NAVY, size=9)
        add_text(slide, title, x + 0.17, 2.91, width - 0.34, 0.46, size=16, fill=WHITE, bold=True)
        add_text(slide, body, x + 0.17, 3.56, width - 0.34, 0.66, size=10.5, fill=MUTED)
        if idx < len(stages) - 1:
            add_line(slide, x + width, 3.34, x + width + gap, 3.34, line=accent, width=2.0)

    add_rect(slide, 0.72, 5.08, 11.90, 1.25, fill=NAVY_2, line=LINE)
    facts = [
        ("9", "CTC epochs", AQUA),
        ("18", "scorer epochs", BLUE),
        ("0", "alignment fallbacks", GREEN),
        ("42", "random seed", GOLD),
    ]
    for idx, (value, label, accent) in enumerate(facts):
        x = 1.00 + idx * 2.88
        add_text(slide, value, x, 5.28, 0.60, 0.45, size=25, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.72, 5.35, 1.58, 0.28, size=10, fill=WHITE, bold=True)
        if idx < 3:
            add_line(slide, x + 2.49, 5.35, x + 2.49, 6.04, line=LINE, width=1.0)
    add_text(slide, "Training discipline: tune epochs on train-only development data, then restart and fit all 2,799 training utterances.", 0.98, 6.52, 11.30, 0.28, size=11, fill=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: submission/WRITEUP.md; experiments/E01-production-model/README.md")
    return slide


def slide_05(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "The objective and metrics match the ordering", 5)
    add_rect(slide, 0.72, 1.93, 5.75, 4.75, fill=NAVY_2, line=LINE)
    add_text(slide, "ORDINAL REGRESSION", 1.03, 2.19, 2.2, 0.22, size=9, fill=AQUA, bold=True)
    add_text(slide, "Two ordered questions", 1.03, 2.58, 3.60, 0.38, size=22, fill=WHITE, bold=True)
    add_text(slide, "q₁ = P(Y ≥ 1)", 1.10, 3.29, 2.20, 0.40, size=19, fill=BLUE, bold=True, font=FONT_MONO)
    add_text(slide, "q₂ = P(Y ≥ 2)", 3.72, 3.29, 2.20, 0.40, size=19, fill=AQUA, bold=True, font=FONT_MONO)
    add_line(slide, 1.34, 4.22, 5.82, 4.22, line=LINE, width=3.0)
    points = [(1.34, "0", "heavy", CORAL), (3.58, "1", "understandable", GOLD), (5.82, "2", "native-like", AQUA)]
    for x, label, desc, accent in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x - 0.12), inch(4.10), inch(0.24), inch(0.24))
        dot.fill.solid(); dot.fill.fore_color.rgb = color(accent); dot.line.fill.background()
        add_text(slide, label, x - 0.23, 4.48, 0.46, 0.30, size=15, fill=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x - 0.62, 4.84, 1.24, 0.24, size=8, fill=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.04, 5.44, 5.10, 0.78, fill=CARD, line=AQUA)
    add_text(slide, "score = 50 × (q₁ + q₂)", 1.30, 5.67, 4.58, 0.30, size=18, fill=WHITE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    add_text(slide, "EVALUATION", 6.92, 2.19, 1.8, 0.22, size=9, fill=BLUE, bold=True)
    add_text(slide, "Balanced MAE", 6.92, 2.58, 3.00, 0.40, size=22, fill=WHITE, bold=True)
    add_text(slide, "Primary metric", 10.67, 2.64, 1.18, 0.23, size=9, fill=AQUA, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "Compute MAE within each true label, then average the three class errors equally.", 6.92, 3.11, 5.13, 0.65, size=13, fill=MUTED)
    add_dot_bullet(slide, "Preserves continuous error magnitude", 6.92, 3.91, 4.9, accent=AQUA, size=12)
    add_dot_bullet(slide, "Stops the 80% majority class dominating", 6.92, 4.44, 4.9, accent=AQUA, size=12)
    add_text(slide, "Secondary checks", 6.92, 5.07, 2.3, 0.30, size=12, fill=WHITE, bold=True)
    add_pill(slide, "QWK", 6.92, 5.49, 0.76, text_color=BLUE)
    add_pill(slide, "MACRO-F1", 7.82, 5.49, 1.18, text_color=GOLD)
    add_pill(slide, "BALANCED ACC.", 9.14, 5.49, 1.54, text_color=CORAL)
    add_pill(slide, "SPEARMAN", 10.82, 5.49, 1.16, text_color=AQUA)
    add_rect(slide, 6.92, 6.08, 5.06, 0.55, fill=CARD_2, line=LINE)
    add_text(slide, "Loss weighting: inverse √ class frequency per phone token", 7.13, 6.24, 4.65, 0.22, size=10.5, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: submission/WRITEUP.md; submission/accent_score/metrics.py")
    return slide


def slide_06(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Validation: audio adds material signal", 6)
    add_text(slide, "BALANCED MAE • LOWER IS BETTER", 0.82, 1.89, 4.25, 0.24, size=9, fill=MUTED, bold=True)
    labels = [
        ("Always 100", 50.00, MUTED),
        ("Strongest static", 32.41, CORAL),
        ("Sequence only", 32.03, GOLD),
        ("Selected acoustic", 22.57, AQUA),
    ]
    chart_x, chart_y, max_w = 2.18, 2.45, 4.10
    for idx, (label, value, accent) in enumerate(labels):
        y = chart_y + idx * 0.78
        add_text(slide, label, 0.80, y + 0.11, 1.24, 0.23, size=10, fill=WHITE, bold=(idx == 3), align=PP_ALIGN.RIGHT)
        add_rect(slide, chart_x, y, max_w, 0.46, fill=CARD, radius=True)
        add_rect(slide, chart_x, y, max_w * value / 50.0, 0.46, fill=accent, radius=True)
        add_text(slide, f"{value:.2f}", chart_x + max_w + 0.17, y + 0.07, 0.70, 0.25, size=11, fill=accent, bold=True)
    add_rect(slide, 0.80, 5.83, 5.80, 0.66, fill=CARD_2, line=AQUA)
    add_text(slide, "−9.454", 1.03, 5.98, 1.04, 0.31, size=20, fill=AQUA, bold=True)
    add_text(slide, "balanced-MAE points vs sequence-only", 2.16, 6.00, 2.74, 0.25, size=10, fill=WHITE, bold=True)
    add_text(slide, "95% CI [−10.556, −8.282]", 4.52, 6.02, 1.82, 0.22, size=8, fill=MUTED, align=PP_ALIGN.RIGHT)

    metric_card(slide, "17.92", "OVERALL MAE", 7.05, 1.98, 2.36, accent=BLUE)
    metric_card(slide, "0.584", "QUADRATIC-WEIGHTED KAPPA", 9.68, 1.98, 2.36, accent=AQUA)
    metric_card(slide, "0.565", "MACRO-F1", 7.05, 3.36, 2.36, accent=GOLD)
    metric_card(slide, "0.551", "SPEARMAN", 9.68, 3.36, 2.36, accent=CORAL)
    add_rect(slide, 7.05, 4.78, 4.99, 1.25, fill=NAVY_2, line=LINE)
    add_text(slide, "2,996 validation phones", 7.30, 5.02, 2.50, 0.31, size=14, fill=WHITE, bold=True)
    add_text(slide, "10,000-draw utterance-bootstrap CI: 21.42–23.80", 7.30, 5.50, 4.42, 0.22, size=9.5, fill=MUTED)
    add_rect(slide, 7.05, 6.22, 4.99, 0.48, fill=CARD_2, line=CORAL)
    add_text(slide, "Important: this split is overlap-heavy—not new-speaker evidence.", 7.25, 6.34, 4.60, 0.21, size=9.5, fill=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: submission/model/metrics.json; submission/WRITEUP.md")
    return slide


def experiment_card(slide, x, y, w, title, result, detail, accent, status="REJECTED"):
    add_rect(slide, x, y, w, 3.35, fill=CARD, line=LINE)
    add_pill(slide, status, x + 0.20, y + 0.22, 0.98, fill=accent, text_color=NAVY, size=8)
    add_text(slide, title, x + 0.20, y + 0.80, w - 0.40, 0.56, size=17, fill=WHITE, bold=True)
    add_text(slide, result, x + 0.20, y + 1.54, w - 0.40, 0.48, size=20, fill=accent, bold=True)
    add_text(slide, detail, x + 0.20, y + 2.18, w - 0.40, 0.82, size=10.5, fill=MUTED)


def slide_07(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "What I tried—and deliberately did not ship", 7)
    experiment_card(
        slide, 0.72, 1.93, 3.73,
        "Bigger encoder", "25.60 vs 22.57",
        "Whisper-small worsened balanced MAE; no small-specific hyperparameter retuning was run.",
        BLUE,
    )
    experiment_card(
        slide, 4.80, 1.93, 3.73,
        "Auxiliary labels", "+0.039 BM delta",
        "Severity and cluster targets produced no reliable gain; class-0 MAE worsened.",
        GOLD,
    )
    experiment_card(
        slide, 8.88, 1.93, 3.73,
        "Full inverse weights", "BM −2.02 • MAE +3.86",
        "Rare-class recall improved, but QWK, macro-F1, majority recall, and calibration regressed.",
        CORAL,
    )
    add_rect(slide, 0.72, 5.64, 11.89, 0.92, fill=CARD_2, line=AQUA, line_width=1.3)
    add_pill(slide, "SELECTED", 0.98, 5.93, 0.94, fill=AQUA, text_color=NAVY, size=8)
    add_text(slide, "Frozen Whisper-tiny + constrained CTC + inverse-√ weighted ordinal head", 2.20, 5.83, 8.56, 0.32, size=15, fill=WHITE, bold=True)
    add_text(slide, "The simplest controlled candidate had the best overall trade-off.", 2.20, 6.22, 8.56, 0.21, size=9.5, fill=MUTED)
    add_text(slide, "Each result comes from its documented protocol; this is not a single shared leaderboard.", 0.72, 6.75, 11.89, 0.20, size=8, fill=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: experiments/E02-whisper-small; E05-auxiliary-labels; E06-scorer-objectives")
    return slide


def slide_08(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "The largest risk is the validation split", 8)
    add_text(slide, "PSEUDO-SPEAKER OVERLAP", 0.83, 1.92, 2.55, 0.23, size=9, fill=MUTED, bold=True)
    add_outline_oval(slide, 0.98, 2.40, 3.28, 3.28, line=BLUE, line_width=2.4)
    add_outline_oval(slide, 3.00, 2.40, 3.28, 3.28, line=AQUA, line_width=2.4)
    add_text(slide, "train", 1.43, 2.74, 1.22, 0.31, size=14, fill=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "validation", 4.62, 2.74, 1.22, 0.31, size=14, fill=AQUA, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.35, 3.47, 2.50, 1.28, fill=CARD_2, line=CORAL, line_width=1.5)
    add_text(slide, "97 / 100", 2.54, 3.66, 2.12, 0.43, size=27, fill=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "validation recordings overlap", 2.54, 4.18, 2.12, 0.24, size=8.5, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "98.0% of validation phones", 1.73, 5.96, 3.80, 0.30, size=14, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "sit inside those inferred voice clusters", 1.73, 6.30, 3.80, 0.24, size=9, fill=MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, 7.04, 1.93, 5.17, 4.75, fill=NAVY_2, line=LINE)
    add_text(slide, "What this changes", 7.38, 2.27, 3.30, 0.40, size=22, fill=WHITE, bold=True)
    add_dot_bullet(slide, "92 / 100 validation prompts also appear in training.", 7.40, 3.03, 4.42, accent=GOLD, size=13, h=0.62)
    add_dot_bullet(slide, "Reported metrics describe a partly seen-speaker / seen-prompt distribution.", 7.40, 3.86, 4.42, accent=CORAL, size=13, h=0.78)
    add_dot_bullet(slide, "A speaker-disjoint benchmark is required before making generalization claims.", 7.40, 4.92, 4.42, accent=AQUA, size=13, h=0.78)
    add_rect(slide, 7.38, 5.99, 4.48, 0.40, fill=CARD, line=LINE)
    add_text(slide, "WavLM clusters are inferred voices—not verified identities.", 7.56, 6.09, 4.12, 0.18, size=8.5, fill=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: data/speaker_clusters/report.md; experiments/E03-speaker-leakage/README.md")
    return slide


def slide_09(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Sniff test: right direction, weak separation", 9)
    add_rect(slide, 0.72, 1.94, 5.76, 4.72, fill=NAVY_2, line=LINE)
    add_text(slide, "CONTROLLED OWN-VOICE PAIR", 1.03, 2.19, 2.85, 0.22, size=9, fill=MUTED, bold=True)
    add_text(slide, "+3.13", 4.90, 2.13, 1.18, 0.42, size=24, fill=AQUA, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "mean points", 4.93, 2.55, 1.15, 0.19, size=8, fill=MUTED, align=PP_ALIGN.RIGHT)
    bars = [("Best American", 70.05, AQUA), ("Non-native", 66.92, CORAL)]
    for idx, (label, value, accent) in enumerate(bars):
        y = 3.12 + idx * 1.05
        add_text(slide, label, 1.04, y + 0.10, 1.37, 0.26, size=11, fill=WHITE, bold=True)
        add_rect(slide, 2.56, y, 3.17, 0.50, fill=CARD, radius=True)
        add_rect(slide, 2.56, y, 3.17 * value / 100.0, 0.50, fill=accent, radius=True)
        add_text(slide, f"{value:.2f}", 5.00, y + 0.08, 0.58, 0.24, size=11, fill=accent, bold=True, align=PP_ALIGN.RIGHT)
    add_dot_bullet(slide, "Only 10 / 20 phones moved higher", 1.04, 5.41, 4.95, accent=GOLD, size=11)
    add_dot_bullet(slide, "Pace differed: 2.76s vs 4.08s", 1.04, 5.91, 4.95, accent=CORAL, size=11)

    add_rect(slide, 6.82, 1.94, 5.40, 4.72, fill=CARD, line=LINE)
    add_text(slide, "HEAVILY ACCENTED PHONES", 7.15, 2.19, 3.25, 0.22, size=9, fill=MUTED, bold=True)
    metric_card(slide, "43.53%", "CORRECTLY SCORED BELOW 25", 7.15, 2.77, 2.15, accent=AQUA)
    metric_card(slide, "11.94%", "INCORRECTLY SCORED ≥ 75", 9.61, 2.77, 2.15, accent=CORAL)
    add_text(slide, "Failure pattern", 7.15, 4.23, 2.15, 0.27, size=12, fill=WHITE, bold=True)
    add_text(slide, "The model reacts when accent also hurts expected-phone recognition, but can miss subtler non-American realizations that Whisper still recognizes confidently.", 7.15, 4.64, 4.62, 1.05, size=12, fill=MUTED)
    add_rect(slide, 7.15, 5.83, 4.62, 0.47, fill=CARD_2, line=CORAL)
    add_text(slide, "Conclusion: marginal utterance-level pass—not a phone-level pass.", 7.35, 5.95, 4.22, 0.21, size=9.5, fill=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: experiments/E07-sniff-tests/SNIFF_TEST.md; E07 and E08")
    return slide


def judge_card(slide, x, title, headline, detail, accent):
    add_rect(slide, x, 2.13, 3.58, 3.66, fill=CARD, line=LINE)
    add_text(slide, title, x + 0.22, 2.41, 3.10, 0.31, size=15, fill=WHITE, bold=True)
    add_text(slide, headline, x + 0.22, 3.06, 3.10, 0.61, size=20, fill=accent, bold=True)
    add_text(slide, detail, x + 0.22, 3.91, 3.10, 1.13, size=11, fill=MUTED)
    add_pill(slide, "NOT SAFE TO RELABEL", x + 0.22, 5.20, 1.78, fill=accent, text_color=NAVY, size=7.5)


def slide_10(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Automated judges did not validate the labels", 10)
    judge_card(slide, 0.72, "GOPT teacher", "5,863 / 5,894 → label 2", "Useful ranking signal, but cross-corpus calibration collapsed near native-like; macro-F1 0.299.", BLUE)
    judge_card(slide, 4.88, "Local audio LLM", "347 / 347 → label 2", "Gemma 4 12B passed structure checks but failed the informativeness gate completely.", GOLD)
    judge_card(slide, 9.04, "GPT audio judge", "40% agreement • 0 label-0", "Balanced 30-item audit: macro-F1 0.299, QWK 0.143; it remained too lenient.", CORAL)
    add_rect(slide, 0.72, 6.15, 11.90, 0.55, fill=CARD_2, line=AQUA)
    add_text(slide, "Next valid step", 0.94, 6.30, 1.32, 0.20, size=9, fill=AQUA, bold=True)
    add_text(slide, "Blinded expert review. The deterministic packets exist; no human-rating ledger exists yet.", 2.41, 6.24, 9.85, 0.28, size=12, fill=WHITE, bold=True)
    add_footer(slide, "Sources: experiments/E09–E13; E11/GOPT_PILOT_RESULTS.md")
    return slide


def slide_11(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Demo: turn a score into a practice loop", 11)
    add_rect(slide, 0.72, 1.88, 7.42, 4.95, fill="F3F7FA", line=LINE)
    add_rect(slide, 0.72, 1.88, 7.42, 0.52, fill="DCE8F0", line=None, radius=False)
    for idx, c in enumerate((CORAL, GOLD, GREEN)):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(0.95 + idx * 0.25), inch(2.04), inch(0.10), inch(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = color(c); dot.line.fill.background()
    add_text(slide, "Pronunciation practice", 1.67, 2.01, 2.80, 0.20, size=9, fill=INK, bold=True)
    add_text(slide, "Sentence to practice", 1.08, 2.71, 2.10, 0.23, size=9, fill="52677B", bold=True)
    add_rect(slide, 1.08, 3.07, 5.72, 0.56, fill=WHITE, line="C8D7E2")
    add_text(slide, "We are both children together.", 1.30, 3.21, 4.85, 0.25, size=12, fill=INK, bold=True)
    add_rect(slide, 6.30, 3.16, 0.31, 0.31, fill=BLUE, radius=True)
    add_text(slide, "▶", 6.30, 3.20, 0.31, 0.18, size=8, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Generated phones • editable", 1.08, 3.94, 2.38, 0.22, size=9, fill="52677B", bold=True)
    phones = ["w", "i", "j", "ɝ", "b", "oʊ", "θ", "tʃ", "ɪ", "l"]
    for idx, phone in enumerate(phones):
        x = 1.08 + idx * 0.54
        add_rect(slide, x, 4.29, 0.43, 0.43, fill="E4EEF4", line="C8D7E2")
        add_text(slide, phone, x, 4.39, 0.43, 0.20, size=10, fill=INK, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.08, 5.08, 2.58, 0.69, fill=BLUE, line=None)
    add_text(slide, "Record / upload audio", 1.30, 5.29, 2.14, 0.24, size=11, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.92, 5.08, 2.88, 0.69, fill=AQUA, line=None)
    add_text(slide, "Score pronunciation", 4.14, 5.29, 2.44, 0.24, size=11, fill=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Score bands", 1.08, 6.10, 1.15, 0.20, size=8, fill="52677B", bold=True)
    for idx, (label, c) in enumerate((("practice", CORAL), ("developing", GOLD), ("strong", GREEN))):
        x = 2.08 + idx * 1.45
        add_rect(slide, x, 6.05, 0.18, 0.18, fill=c, radius=True)
        add_text(slide, label, x + 0.25, 6.03, 0.94, 0.19, size=7.5, fill="52677B")

    add_text(slide, "Product choices", 8.57, 1.98, 3.10, 0.35, size=22, fill=WHITE, bold=True)
    add_dot_bullet(slide, "Microphone or upload", 8.58, 2.72, 3.38, accent=AQUA, size=13)
    add_dot_bullet(slide, "Sentence playback", 8.58, 3.25, 3.38, accent=BLUE, size=13)
    add_dot_bullet(slide, "Automatic, editable phonemes", 8.58, 3.78, 3.38, accent=GOLD, size=13)
    add_dot_bullet(slide, "Ordered per-phone feedback", 8.58, 4.31, 3.38, accent=CORAL, size=13)
    add_rect(slide, 8.58, 5.06, 3.68, 1.08, fill=CARD_2, line=AQUA)
    add_text(slide, "Difficulty ≠ model score", 8.80, 5.25, 3.24, 0.28, size=13, fill=AQUA, bold=True)
    add_text(slide, "Keep raw scores fixed; adapt only the coaching threshold per phone and learner.", 8.80, 5.65, 3.24, 0.36, size=9.5, fill=WHITE)
    add_text(slide, "Run locally:  cd submission  •  uv run python demo_app.py", 8.58, 6.48, 3.68, 0.20, size=7.5, fill=MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: submission/demo_app.py; submission/README.md; submission/WRITEUP.md")
    return slide


def slide_12(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "Takeaway: useful signal, honest uncertainty", 12)
    add_rich_text(
        slide,
        [
            ("The model learned ", {"fill": WHITE, "bold": True}),
            ("useful acoustic signal", {"fill": AQUA, "bold": True}),
            (".\nThe evaluation has not yet proven ", {"fill": WHITE, "bold": True}),
            ("new-speaker robustness", {"fill": CORAL, "bold": True}),
            (".", {"fill": WHITE, "bold": True}),
        ],
        0.78,
        1.95,
        11.72,
        1.18,
        size=28,
    )
    steps = [
        ("01", "Benchmark", "Collect a genuinely speaker-disjoint, expert-rated evaluation set.", AQUA),
        ("02", "Calibrate", "Review label disagreements and fit phone-specific calibration with uncertainty.", BLUE),
        ("03", "Improve", "Test a more accent-sensitive encoder across seeds and matched-pace voice pairs.", GOLD),
    ]
    for idx, (num, title, body, accent) in enumerate(steps):
        x = 0.78 + idx * 4.02
        add_rect(slide, x, 3.47, 3.63, 2.03, fill=CARD, line=LINE)
        add_pill(slide, num, x + 0.22, 3.69, 0.50, fill=accent, text_color=NAVY, size=9)
        add_text(slide, title, x + 0.22, 4.22, 2.98, 0.34, size=18, fill=WHITE, bold=True)
        add_text(slide, body, x + 0.22, 4.78, 3.05, 0.49, size=10.5, fill=MUTED)
    add_rect(slide, 0.78, 5.88, 11.67, 0.69, fill=CARD_2, line=CORAL)
    add_text(slide, "Ship today", 1.02, 6.07, 1.12, 0.25, size=10, fill=CORAL, bold=True)
    add_text(slide, "A coaching prototype with visible uncertainty—not a high-stakes judgment of identity, ability, or proficiency.", 2.20, 6.01, 9.91, 0.34, size=12.5, fill=WHITE, bold=True)
    add_text(slide, "Questions  •  discussion  •  live demo", 0.78, 6.75, 11.67, 0.24, size=11, fill=AQUA, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Sources: submission/WRITEUP.md; experiments/README.md")
    return slide


def validate(prs: Presentation) -> None:
    if len(prs.slides) != 12:
        raise RuntimeError(f"Expected 12 slides, found {len(prs.slides)}")
    for index, slide in enumerate(prs.slides, start=1):
        if not slide.shapes:
            raise RuntimeError(f"Slide {index} is empty")
        for shape in slide.shapes:
            # Decorative title-slide circles intentionally bleed off canvas.
            if index == 1 and (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > prs.slide_width
                or shape.top + shape.height > prs.slide_height
            ):
                continue
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {index} has a shape outside the top/left canvas")
            if shape.left + shape.width > prs.slide_width + inch(0.02):
                raise RuntimeError(f"Slide {index} has a shape outside the right canvas")
            if shape.top + shape.height > prs.slide_height + inch(0.02):
                raise RuntimeError(f"Slide {index} has a shape outside the bottom canvas")


def build() -> Path:
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    prs.core_properties.title = "Phone-Level Accentedness Scoring"
    prs.core_properties.subject = "ML take-home challenge presentation"
    prs.core_properties.author = "Aviad Arnias"
    prs.core_properties.keywords = "speech, pronunciation, accentedness, ordinal regression, Whisper"
    prs.core_properties.comments = "Generated from presentation/generate_slides.py"

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    validate(prs)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    output = build()
    print(f"Wrote {output.relative_to(ROOT)}")
